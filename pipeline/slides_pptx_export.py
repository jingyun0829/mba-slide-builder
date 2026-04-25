"""Export the current outline to .pptx for co-teachers / LMS / offline edits.

Reads the same slide-centric outline JSON the HTML renderer uses and
produces a .pptx with rich styling preserved (bold, smaller, red,
monospace for code, syntax-highlighted code via Pygments, homework block).
Image-type slides render as italic placeholder text — SVG doesn't translate
to pptx without extra libraries.
"""
from __future__ import annotations
import json
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches, Emu

_TITLE_HINTS = ["title slide", "title", "cover"]
_CONTENT_HINTS = ["title and content", "content", "bullet", "body"]
RED = RGBColor(0xC0, 0x00, 0x00)

# --- Theme palette for fancy decoration (A + B + D upgrades) ---
TEAL = RGBColor(0x0d, 0x94, 0x88)
TEAL_DARK = RGBColor(0x04, 0x2f, 0x2e)
MINT = RGBColor(0xa7, 0xf3, 0xd0)
WHITE = RGBColor(0xff, 0xff, 0xff)
FOOTER_GRAY = RGBColor(0x94, 0xa3, 0xb8)

# Slide-type emoji icons for title prefix
TYPE_ICONS = {
    "concept":    "💡",
    "question":   "❓",
    "example":    "🏢",
    "image":      "🖼️",
    "code":       "💻",
    "exercise":   "⚙️",
    "discussion": "💬",
    "summary":    "✨",
}

# --- Pygments-based syntax highlighting for code slides ---
# Tomorrow Night theme (dark background, light-on-dark colors) — matches HTML's Prism theme.
_CODE_BG_COLOR = RGBColor(0x2D, 0x2D, 0x2D)  # dark charcoal background
_TOKEN_COLOR_RULES = [
    ("Comment",           RGBColor(0x99, 0x99, 0x99)),  # muted gray, italic
    ("String",            RGBColor(0x99, 0xCC, 0x99)),  # soft green — string literals
    ("Number",            RGBColor(0xF9, 0x91, 0x57)),  # orange — numbers
    ("Keyword.Namespace", RGBColor(0xCC, 0x99, 0xCC)),  # light purple — import, from
    ("Keyword.Constant",  RGBColor(0xF9, 0x91, 0x57)),  # orange — True/False/None
    ("Keyword",           RGBColor(0xCC, 0x99, 0xCC)),  # light purple — def, for, if, return
    ("Operator.Word",     RGBColor(0xCC, 0x99, 0xCC)),  # light purple — and/or/not/in
    ("Name.Function",     RGBColor(0xF2, 0x77, 0x7A)),  # coral — function names
    ("Name.Builtin",      RGBColor(0xFF, 0xCC, 0x66)),  # yellow — print, len, range
    ("Name.Class",        RGBColor(0x66, 0xCC, 0xCC)),  # cyan — class names
    ("Name.Decorator",    RGBColor(0xF2, 0x77, 0x7A)),  # coral — @decorator
]
_DEFAULT_CODE_COLOR = RGBColor(0xCC, 0xCC, 0xCC)  # foreground light gray


def _token_color(token_type):
    name = str(token_type)  # e.g. "Token.Keyword.Namespace"
    for key, color in _TOKEN_COLOR_RULES:
        if key in name:
            return color
    return _DEFAULT_CODE_COLOR


def _is_comment_token(token_type):
    return "Comment" in str(token_type)

def _find_layout(prs, hints, fallback_idx):
    for hint in hints:
        for layout in prs.slide_layouts:
            if hint.lower() in (layout.name or "").lower():
                return layout
    if fallback_idx < len(prs.slide_layouts):
        return prs.slide_layouts[fallback_idx]
    return prs.slide_layouts[0]

def _set_title(slide, text):
    if slide.shapes.title is not None:
        slide.shapes.title.text = text or ""

def _first_body_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            return ph
    return None

def _clear_slides(prs):
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        xml_slides.remove(sld_id)

def _normalize_line(line):
    if isinstance(line, str):
        return line, False, False, False
    if isinstance(line, dict):
        return (str(line.get("text","")), bool(line.get("bold",False)),
                bool(line.get("small",False)), bool(line.get("red",False)))
    return str(line), False, False, False

def _add_styled_line(tf, text, bold, small, red, first=False, base_size=24):
    if not text.strip(): return
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.level = 0
    for run in p.runs:
        run.font.size = Pt(18 if small else base_size)
        if bold or red:
            run.font.bold = True
        if red:
            run.font.color.rgb = RED

# ============================================================
# Fancy decoration helpers (A: title accent + icon, B: callout takeaway, D: footer)
# ============================================================

def _prepend_type_icon(slide, slide_type: str) -> None:
    """Prefix the title text with a type-specific emoji (concept 💡, question ❓, etc.)."""
    icon = TYPE_ICONS.get((slide_type or "").lower(), "")
    if not icon or slide.shapes.title is None:
        return
    cur = slide.shapes.title.text or ""
    if not cur or cur.startswith(icon):
        return
    slide.shapes.title.text = f"{icon}  {cur}"
    # Re-apply title formatting (set_title resets to default)
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        if run.font.size is None:
            run.font.size = Pt(32)


def _add_title_accent_bar(slide, color=TEAL) -> None:
    """Add a short colored bar just under the title placeholder."""
    if slide.shapes.title is None:
        return
    t = slide.shapes.title
    try:
        left = t.left
        top = t.top + t.height - Inches(0.05)
        width = Inches(2.0)
        height = Inches(0.08)
    except Exception:
        return
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_callout_takeaway(slide, prs, text: str) -> None:
    """Replace the inline red takeaway with a rounded teal callout box at the bottom."""
    if not text or not text.strip():
        return
    sw = prs.slide_width
    sh = prs.slide_height
    margin = Inches(0.5)
    height = Inches(1.05)
    width = sw - 2 * margin
    left = margin
    top = sh - margin - height - Inches(0.35)  # leave room for footer

    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    callout.fill.solid()
    callout.fill.fore_color.rgb = TEAL
    callout.line.fill.background()
    # Soft shadow effect (subtle)
    try:
        callout.shadow.inherit = False
    except Exception:
        pass

    tf = callout.text_frame
    tf.margin_left = Inches(0.30)
    tf.margin_right = Inches(0.30)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.16)
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "▌  KEY INSIGHT"
    for run in p1.runs:
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = MINT
        run.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = text.strip()
    for run in p2.runs:
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"


def _add_footer(slide, prs, page_num: int, total_pages: int,
                left_text: str = "") -> None:
    """Add small footer at bottom: 'left_text' on left, 'N / Total' on right."""
    sw = prs.slide_width
    sh = prs.slide_height
    margin = Inches(0.4)
    footer_h = Inches(0.25)
    top = sh - footer_h - Inches(0.08)

    half_w = (sw - 2 * margin) / 2

    # Left footer
    if left_text:
        if len(left_text) > 65:
            left_text = left_text[:62] + "…"
        tb_l = slide.shapes.add_textbox(margin, top, half_w, footer_h)
        ftl = tb_l.text_frame
        ftl.margin_left = ftl.margin_right = Emu(0)
        ftl.margin_top = ftl.margin_bottom = Emu(0)
        p = ftl.paragraphs[0]
        p.text = left_text
        for run in p.runs:
            run.font.size = Pt(8)
            run.font.color.rgb = FOOTER_GRAY
            run.font.name = "Calibri"

    # Right footer (page N / Total)
    tb_r = slide.shapes.add_textbox(margin + half_w, top, half_w, footer_h)
    ftr = tb_r.text_frame
    ftr.margin_left = ftr.margin_right = Emu(0)
    ftr.margin_top = ftr.margin_bottom = Emu(0)
    p = ftr.paragraphs[0]
    p.text = f"{page_num} / {total_pages}"
    p.alignment = PP_ALIGN.RIGHT
    for run in p.runs:
        run.font.size = Pt(8)
        run.font.color.rgb = FOOTER_GRAY
        run.font.name = "Calibri"


def _render_title_slide(prs, title, duration):
    layout = _find_layout(prs, _TITLE_HINTS, fallback_idx=0)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    body = _first_body_placeholder(slide)
    if body is not None:
        body.text_frame.text = f"{duration} min"

def _render_content_slide(prs, title, lines, key_takeaway="", slide_type="concept"):
    layout = _find_layout(prs, _CONTENT_HINTS, fallback_idx=1)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    body = _first_body_placeholder(slide)
    if body is None: return slide
    tf = body.text_frame
    tf.clear()
    first = True
    for line in (lines or []):
        text, bold, small, red = _normalize_line(line)
        if not text.strip(): continue
        _add_styled_line(tf, text, bold, small, red, first=first)
        first = False

    # === Fancy upgrades A + B ===
    _prepend_type_icon(slide, slide_type)
    _add_title_accent_bar(slide, color=TEAL)
    if key_takeaway and key_takeaway.strip():
        _add_callout_takeaway(slide, prs, key_takeaway)
    return slide

def _render_image_slide(prs, title, image_hint):
    layout = _find_layout(prs, _CONTENT_HINTS, fallback_idx=1)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    body = _first_body_placeholder(slide)
    if body is None: return
    tf = body.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = "[ INSERT DIAGRAM / IMAGE HERE ]"
    for run in p1.runs:
        run.font.size = Pt(20)
        run.font.italic = True
    if image_hint:
        p2 = tf.add_paragraph()
        p2.text = image_hint
        for run in p2.runs:
            run.font.size = Pt(14)
            run.font.italic = True
    # Fancy upgrades on image slides
    _prepend_type_icon(slide, "image")
    _add_title_accent_bar(slide, color=TEAL)
    return slide

def _render_code_step_slide(prs, step_title, code, language, explanation=""):
    layout = _find_layout(prs, _CONTENT_HINTS, fallback_idx=1)
    slide = prs.slides.add_slide(layout)
    lang_tag = f" ({language})" if language else ""
    _set_title(slide, f"{step_title}{lang_tag}")
    body = _first_body_placeholder(slide)
    if body is None: return
    tf = body.text_frame
    tf.clear()

    # Dark charcoal background for the body placeholder (matches HTML Tomorrow theme).
    try:
        body.fill.solid()
        body.fill.fore_color.rgb = _CODE_BG_COLOR
    except Exception:
        pass  # some templates / shape types don't support fill overrides

    # Fancy upgrades on code slides
    _prepend_type_icon(slide, "code")
    _add_title_accent_bar(slide, color=TEAL)

    # Try Pygments-based syntax highlighting first; fall back to plain monospace.
    try:
        from pygments import lex
        from pygments.lexers import get_lexer_by_name
        try:
            lexer = get_lexer_by_name(language or "python")
        except Exception:
            lexer = get_lexer_by_name("python")
        tokens = list(lex(code or "", lexer))
        _render_tokens_into(tf, tokens)
        if explanation and explanation.strip():
            _append_explanation(tf, explanation.strip())
        return slide
    except ImportError:
        pass  # pygments not installed — use plain fallback
    except Exception:
        pass  # any lexing error — plain fallback

    # Plain monospace fallback (no colors)
    code_lines = (code or "").split("\n") or [""]
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line else " "
        p.level = 0
        for run in p.runs:
            run.font.name = "Consolas"
            run.font.size = Pt(16)


def _render_tokens_into(tf, tokens):
    """Render a list of (TokenType, text) tokens into the text_frame with colored runs."""
    # Split tokens by newlines into lines so we can create one paragraph per code line.
    lines_of_tokens: list[list[tuple]] = [[]]
    for token_type, value in tokens:
        if not value:
            continue
        parts = value.split("\n")
        for i, part in enumerate(parts):
            if part:
                lines_of_tokens[-1].append((token_type, part))
            if i < len(parts) - 1:
                lines_of_tokens.append([])

    for line_idx, line_tokens in enumerate(lines_of_tokens):
        p = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
        p.level = 0
        if not line_tokens:
            # empty line — add a space run to preserve vertical rhythm
            r = p.add_run()
            r.text = " "
            r.font.name = "Consolas"
            r.font.size = Pt(15)
            continue
        for token_type, text in line_tokens:
            r = p.add_run()
            r.text = text
            r.font.name = "Consolas"
            r.font.size = Pt(15)
            r.font.color.rgb = _token_color(token_type)
            if _is_comment_token(token_type):
                r.font.italic = True


def _append_explanation(tf, explanation: str):
    """Append a blank-line + italic light-gray explanation below the code block."""
    # Spacer
    p_sp = tf.add_paragraph()
    p_sp.text = " "
    for r in p_sp.runs:
        r.font.size = Pt(8)
    # Explanation
    p_ex = tf.add_paragraph()
    p_ex.text = explanation
    for r in p_ex.runs:
        r.font.name = "Calibri"
        r.font.size = Pt(14)
        r.font.italic = True
        r.font.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)  # light gray on dark bg


def _render_code(prs, outer_title, code_obj):
    if not code_obj: return
    lang = code_obj.get("language","")
    steps = code_obj.get("steps") or []
    if not steps:
        raw = code_obj.get("content","")
        if raw:
            _render_code_step_slide(prs, outer_title or "Code", raw, lang,
                                    explanation=code_obj.get("caption", ""))
        return
    for idx, s in enumerate(steps, 1):
        raw_desc = ((s or {}).get("description") or "").strip().rstrip("—-: ").strip()
        desc = raw_desc if raw_desc else f"Step {idx}"
        if not desc.lower().startswith("step"):
            desc = f"Step {idx}: {desc}"
        code_text = ((s or {}).get("code") or "").strip()
        if not code_text:
            continue
        explanation = ((s or {}).get("explanation") or "").strip()
        _render_code_step_slide(prs, desc, code_text, lang, explanation=explanation)

def _render_activity(prs, activity):
    if not activity:
        return
    title = activity.get("title", "Warm-up Activity")
    atype = activity.get("type", "web_link")
    duration = activity.get("duration_minutes", 10)

    # Section divider
    _render_content_slide(prs, "Warm-up Activity", [])

    # Intro
    intro_lines = []
    if activity.get("scenario"):
        intro_lines.append(activity["scenario"])
    if activity.get("learning_goal"):
        intro_lines.append({"text": f"Goal: {activity['learning_goal']}", "small": True})
    _render_content_slide(prs, f"{title}  ({duration} min)", intro_lines)

    # Instructions
    instr_lines = []
    if activity.get("instructions"):
        instr_lines.append(activity["instructions"])
    if atype == "excel_simulation":
        instr_lines.append({"text": "Open the Excel file provided with this deck.", "bold": True})
    elif atype == "web_link":
        if activity.get("url"):
            instr_lines.append({"text": f"Visit: {activity['url']}", "bold": True})
        if activity.get("source_name"):
            instr_lines.append({"text": f"Source: {activity['source_name']}", "small": True})
        if activity.get("what_to_do"):
            instr_lines.append(activity["what_to_do"])
    _render_content_slide(prs, "What to do", instr_lines)

    # Debrief
    if activity.get("debrief_questions"):
        _render_content_slide(
            prs, "Debrief",
            activity["debrief_questions"],
            key_takeaway=(activity.get("facilitation_notes", "").split(".")[0]
                          if activity.get("facilitation_notes") else ""),
        )


def _render_homework(prs, hw):
    _render_content_slide(prs, "Homework", [hw.get("title","Homework")])
    if hw.get("problem_statement"):
        _render_content_slide(prs, "Your Task", [hw["problem_statement"].replace("\n"," ")])
    ds = hw.get("dataset") or {}
    if ds:
        lines = []
        if ds.get("description"): lines.append(ds["description"])
        if ds.get("source"): lines.append({"text": f"Source: {ds['source']}", "small": True})
        if ds.get("columns"):
            cols = ds["columns"]
            lines.append({"text": f"Columns: {', '.join(cols[:8])}" + (" …" if len(cols) > 8 else ""), "small": True})
        if lines: _render_content_slide(prs, "Dataset", lines)
    if hw.get("deliverables"): _render_content_slide(prs, "Deliverables", hw["deliverables"])
    if hw.get("hints"): _render_content_slide(prs, "Hints", hw["hints"])
    if hw.get("grading_rubric"): _render_content_slide(prs, "Grading", [hw["grading_rubric"]])

def export_pptx(outline_json, output_path, template_path=None):
    outline = json.loads(outline_json)
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        _clear_slides(prs)
    else:
        prs = Presentation()
    session_title = outline.get("session_title", "Untitled")
    _render_title_slide(prs, session_title, outline.get("duration_minutes", 90))
    slides_list = outline.get("slides") or []
    has_lo = any((s.get("title", "").lower().strip() in ("learning objectives", "objectives")) for s in slides_list)
    if outline.get("learning_objectives") and not has_lo:
        _render_content_slide(prs, "Learning Objectives", outline["learning_objectives"], slide_type="summary")
    for s in slides_list:
        stype = (s.get("type") or "concept").lower().strip()
        title = s.get("title", "")
        lines = s.get("lines", [])
        kt = s.get("key_takeaway", "")
        if stype == "title": continue
        elif stype == "image": _render_image_slide(prs, title, s.get("image_hint", ""))
        elif stype == "code": _render_code(prs, title, s.get("code") or {})
        else: _render_content_slide(prs, title, lines, key_takeaway=kt, slide_type=stype)
    if outline.get("activity"):
        _render_activity(prs, outline["activity"])
    if outline.get("homework"):
        _render_homework(prs, outline["homework"])

    # === D) Footer pass: add session title + page number to every slide except cover (slide 0) ===
    all_slides = list(prs.slides)
    total_pages = len(all_slides)
    for idx, slide in enumerate(all_slides):
        if idx == 0:
            continue  # skip the title cover slide
        try:
            _add_footer(slide, prs, page_num=idx + 1, total_pages=total_pages,
                        left_text=session_title)
        except Exception:
            pass  # don't fail the whole build if a single footer breaks

    prs.save(output_path)
    return output_path
