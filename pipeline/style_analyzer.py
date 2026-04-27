"""Extract a teaching-style profile from the instructor's past .pptx decks.

Reads EVERY shape's text frame — not just placeholders — so instructors who
compose slides with text boxes + images get their style captured accurately.
"""
from __future__ import annotations
import json, os, statistics
from pathlib import Path
from anthropic import Anthropic
from pptx import Presentation

_client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
_MODEL = os.getenv("ANTHROPIC_MODEL", "claude-opus-4-6")
_PROMPT_DIR = Path(__file__).resolve().parent.parent / "prompts"
_PICTURE_SHAPE_TYPE = 13


def _extract_deck(pptx_path):
    prs = Presentation(pptx_path)
    slides = []
    for slide in prs.slides:
        lines = []
        images = 0
        for shape in slide.shapes:
            if shape.shape_type == _PICTURE_SHAPE_TYPE:
                images += 1
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
        slides.append({"lines": lines, "images": images})
    return {"name": Path(pptx_path).name, "slides": slides}


def _compute_quant(decks):
    all_slides = [s for d in decks for s in d["slides"]]
    if not all_slides:
        return {"decks_analyzed": len(decks), "total_slides_analyzed": 0}
    slide_counts = [len(d["slides"]) for d in decks]
    lines_per_slide = [len(s["lines"]) for s in all_slides]
    titles = [s["lines"][0] for s in all_slides if s["lines"]]
    body_lines = [l for s in all_slides for l in s["lines"][1:]]
    title_words = [len(t.split()) for t in titles]
    body_words = [len(l.split()) for l in body_lines]
    question_titles = sum(1 for t in titles if "?" in t)
    short_titles = sum(1 for t in titles if len(t.split()) <= 5)
    slides_with_images = sum(1 for s in all_slides if s["images"] > 0)

    def _avg(xs):
        return round(statistics.mean(xs), 2) if xs else 0
    def _med(xs):
        return round(statistics.median(xs), 1) if xs else 0

    return {
        "decks_analyzed": len(decks),
        "total_slides_analyzed": sum(slide_counts),
        "avg_slides_per_lecture": _avg(slide_counts),
        "avg_lines_per_slide": _avg(lines_per_slide),
        "median_lines_per_slide": _med(lines_per_slide),
        "avg_words_per_title": _avg(title_words),
        "avg_words_per_body_line": _avg(body_words),
        "question_title_ratio": round(question_titles/len(titles),2) if titles else 0,
        "short_title_ratio": round(short_titles/len(titles),2) if titles else 0,
        "image_slide_ratio": round(slides_with_images/len(all_slides),2),
    }


def _build_style_sample(decks, max_slides=80):
    chunks = []
    emitted = 0
    for d_idx, d in enumerate(decks, 1):
        chunks.append(f"=== Deck {d_idx}: {d['name']} ({len(d['slides'])} slides) ===")
        for s_idx, s in enumerate(d["slides"], 1):
            if emitted >= max_slides:
                chunks.append(f"(...and {len(d['slides']) - s_idx + 1} more slides)")
                break
            img_tag = f" [{s['images']} image{'s' if s['images']>1 else ''}]" if s['images'] else ""
            chunks.append(f"Slide {s_idx}{img_tag}")
            for line in s["lines"]:
                chunks.append(f"  {line}")
            chunks.append("")
            emitted += 1
        if emitted >= max_slides:
            break
    return "\n".join(chunks)


def _robust_json_parse(text: str) -> dict:
    """Parse JSON from a potentially messy LLM response.

    Handles all of these failure modes seen in production:
      1. Plain JSON                                          → json.loads
      2. ```json ... ``` markdown fences                     → strip fences
      3. ``` ... ``` (no language)                           → strip fences
      4. Preamble like 'Here is your analysis:\n{...}'       → find first '{'
      5. Trailing notes after the JSON                       → find last '}'
      6. Multiple JSON objects (use the largest balanced)    → bracket counting
      7. Smart quotes / Unicode quotes                       → normalize
      8. Trailing commas (a Claude habit)                    → strip them
    """
    import re

    if not text:
        raise ValueError("Empty response from style-extraction model.")

    s = text.strip()

    # Strip markdown fences if present
    fence_match = re.search(r'```(?:json|JSON)?\s*\n?(.*?)\n?```', s, re.DOTALL)
    if fence_match:
        s = fence_match.group(1).strip()

    # Try a direct parse first
    try:
        return json.loads(s)
    except json.JSONDecodeError:
        pass

    # Find the first '{' and the matching last '}'
    start = s.find('{')
    end = s.rfind('}')
    if start == -1 or end == -1 or end < start:
        raise ValueError(f"Could not locate a JSON object in the response. "
                         f"First 200 chars: {s[:200]!r}")
    candidate = s[start:end + 1]

    # Normalize smart quotes that Claude sometimes emits
    candidate = (candidate
                 .replace('\u201c', '"').replace('\u201d', '"')
                 .replace('\u2018', "'").replace('\u2019', "'"))
    # Strip trailing commas inside objects/arrays
    candidate = re.sub(r',(\s*[}\]])', r'\1', candidate)

    try:
        return json.loads(candidate)
    except json.JSONDecodeError as e:
        # One last attempt: walk forward and find the largest balanced { ... }
        depth = 0
        in_string = False
        escape = False
        last_balanced_end = None
        for i, ch in enumerate(candidate):
            if escape:
                escape = False
                continue
            if ch == '\\':
                escape = True
                continue
            if ch == '"':
                in_string = not in_string
                continue
            if in_string:
                continue
            if ch == '{':
                depth += 1
            elif ch == '}':
                depth -= 1
                if depth == 0:
                    last_balanced_end = i
                    break
        if last_balanced_end is not None:
            try:
                return json.loads(candidate[:last_balanced_end + 1])
            except json.JSONDecodeError:
                pass
        raise ValueError(
            f"Style-extraction returned non-JSON content. "
            f"Parse error: {e}. "
            f"First 300 chars of response: {text[:300]!r}"
        )


def _qualitative_from_claude(decks, quant):
    system = (_PROMPT_DIR / "style_extraction.md").read_text()
    sample = _build_style_sample(decks)
    user_msg = f"""Quantitative stats from these decks:
{json.dumps(quant, indent=2)}

Raw slide content (every shape's text, in reading order, across the sample):
{sample}

Return the style-profile JSON matching the schema. JSON only — no preamble, no markdown fences, no commentary."""
    resp = _client.messages.create(model=_MODEL, max_tokens=2500, system=system,
                                    messages=[{"role":"user","content":user_msg}])
    text = resp.content[0].text
    return _robust_json_parse(text)


def extract_style_profile(pptx_paths):
    if not pptx_paths:
        raise ValueError("No .pptx files provided.")
    decks = [_extract_deck(p) for p in pptx_paths]
    quant = _compute_quant(decks)
    qual = _qualitative_from_claude(decks, quant)
    return {"quantitative": quant, "qualitative": qual, "sources": [Path(p).name for p in pptx_paths]}


def save_style_profile(profile, path="style_profiles/current.json"):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    Path(path).write_text(json.dumps(profile, indent=2))
    return path


def load_style_profile(path="style_profiles/current.json"):
    p = Path(path)
    return json.loads(p.read_text()) if p.exists() else None


def _compute_auto_dimensions(profile: dict) -> dict[str, int]:
    """The auto-detected 6 scores derived purely from the uploaded pptx files.
    No user overrides applied. Use this when you want to know 'what the
    teacher's actual decks look like' (vs. their manual preferences)."""
    if not profile:
        return {}
    q = profile.get("quantitative", {}) or {}
    qual = profile.get("qualitative", {}) or {}

    def _clamp(x: float) -> int:
        return max(0, min(100, int(round(x))))

    # ----- Derived from quantitative numbers -----
    awpl = float(q.get("avg_words_per_body_line") or 12)
    concise = _clamp(100 - (awpl - 5) * (100 / 13))

    qtr = float(q.get("question_title_ratio") or 0)
    questions = _clamp(qtr * 100 * 1.4)

    isr = float(q.get("image_slide_ratio") or 0)
    visual = _clamp(isr * 100 * 1.5)

    aspl = float(q.get("avg_slides_per_lecture") or 25)
    pacing = _clamp((aspl / 50) * 100)

    # ----- Derived from qualitative text (keyword scoring) -----
    tone = (qual.get("tone") or "").lower()
    conv_keywords = ("conversational", "informal", "casual", " we ", "you ", "friendly")
    conversational = _clamp(85 if any(k in tone for k in conv_keywords) else 45)

    ex_density = (qual.get("example_density") or "").lower()
    if "high" in ex_density:
        real_cases = 90
    elif "medium" in ex_density:
        real_cases = 60
    elif "low" in ex_density:
        real_cases = 25
    else:
        real_cases = 50

    return {
        "Concise lines":     concise,
        "Question-driven":   questions,
        "Visual-rich":       visual,
        "Rapid pacing":      pacing,
        "Conversational":    conversational,
        "Real-case heavy":   real_cases,
    }


def compute_dimensions(profile: dict, ignore_user_overrides: bool = False) -> dict[str, int]:
    """Return the 6 dimension scores. Honors `profile['user_dimensions']`
    overrides if the teacher has dragged the sliders to customize their style.

    Pass ignore_user_overrides=True to get the raw auto-detected values
    (used to display 'auto' hints next to the sliders)."""
    auto = _compute_auto_dimensions(profile)
    if ignore_user_overrides or not profile:
        return auto
    user_dims = profile.get("user_dimensions") or {}
    if not user_dims:
        return auto
    # Merge: user value wins for any dim it specified, auto wins otherwise.
    return {k: int(user_dims[k]) if k in user_dims else v for k, v in auto.items()}


def _dimension_to_instruction(label: str, value: int) -> str:
    """Convert a 0-100 dimension score into a natural-language style instruction
    the LLM can act on. Returns empty string if the value is moderate (40-60),
    since 'be average' isn't useful guidance."""
    if 40 <= value <= 60:
        return ""

    intensity = "STRONGLY" if value >= 80 or value <= 20 else "MODERATELY"

    if label == "Concise lines":
        if value >= 60:
            return (f"- {intensity} prefer SHORT, punchy bullet lines (5-10 words each). "
                    "Cut filler words. Each bullet = one idea, fragment-style is OK.")
        else:
            return (f"- {intensity} prefer LONGER, descriptive bullet lines (15-25 words). "
                    "Full sentences with context, not telegram fragments.")

    if label == "Question-driven":
        if value >= 60:
            return (f"- {intensity} use QUESTION-form slide titles ('Why does X happen?', "
                    "'What if Y?'). Aim for {pct}% of titles as questions."
                    .format(pct=value))
        else:
            return (f"- {intensity} use DECLARATIVE titles ('Three causes of X', 'How Y works'). "
                    "Avoid posing every title as a question.")

    if label == "Visual-rich":
        if value >= 60:
            return (f"- {intensity} include MANY image-type slides ({value}% of slides should be type='image' "
                    "with diagram-worthy image_hint). Visuals carry the lecture, text supports.")
        else:
            return (f"- {intensity} keep image slides RARE ({value}% max). Text + bullets carry the lecture; "
                    "use images only when truly necessary.")

    if label == "Rapid pacing":
        if value >= 60:
            return (f"- {intensity} use RAPID PACING — many short slides (40+ per 90 min lecture). "
                    "Each slide = one focused idea, advance quickly.")
        else:
            return (f"- {intensity} use METHODICAL PACING — fewer, denser slides (15-22 per 90 min). "
                    "Each slide explores an idea more thoroughly.")

    if label == "Conversational":
        if value >= 60:
            return (f"- {intensity} use CONVERSATIONAL tone — 'we', 'you', 'imagine', informal phrasing. "
                    "Address the student directly. Sound like a friendly mentor, not a textbook.")
        else:
            return (f"- {intensity} use FORMAL ACADEMIC tone — third person, precise vocabulary, "
                    "no 'we'/'you'. Sound like a published textbook.")

    if label == "Real-case heavy":
        if value >= 60:
            return (f"- {intensity} ground every concept in REAL COMPANIES / CASES (Amazon, Tesla, Netflix, "
                    "specific named situations). Avoid generic 'a company', use named examples.")
        else:
            return (f"- {intensity} keep examples ABSTRACT and theoretical. Use 'a firm', 'a manager', "
                    "generic scenarios over named-company cases.")

    return ""


def _user_dimensions_to_prompt_block(dims: dict) -> str:
    """Convert the user's 6 customized dimension scores into prompt instructions.
    These take precedence over the auto-detected style — if the user dragged the
    'Visual-rich' slider to 90 even though their old decks scored 30, the AI
    should generate visual-heavy slides."""
    if not dims:
        return ""
    instructions = []
    for label, value in dims.items():
        instr = _dimension_to_instruction(label, value)
        if instr:
            instructions.append(instr)
    if not instructions:
        return ""
    return (
        "\n\n=== INSTRUCTOR'S STYLE PREFERENCES (override the auto-detected style above) ===\n"
        "The instructor manually adjusted these dials in Stage 2. Honor them OVER what their\n"
        "old decks suggest — they're telling you how they WANT to teach, not how they used to.\n\n"
        + "\n".join(instructions)
    )


def profile_to_prompt_block(profile):
    if not profile:
        return ""
    q = profile.get("quantitative", {}) or {}
    qual = profile.get("qualitative", {}) or {}
    lines = [
        "INSTRUCTOR TEACHING STYLE — match this as closely as possible.",
        "",
        "Quantitative targets (match these averages tightly):",
        f"- Slides per lecture: {q.get('avg_slides_per_lecture','?')} (within plus/minus 15%)",
        f"- Lines per slide: {q.get('avg_lines_per_slide','?')} avg, {q.get('median_lines_per_slide','?')} median",
        f"- Words per title: {q.get('avg_words_per_title','?')} avg",
        f"- Words per body line: {q.get('avg_words_per_body_line','?')} avg",
        f"- Question-form titles: {int(q.get('question_title_ratio',0)*100)}% — use question titles at this rate",
        f"- Short titles (<=5 words): {int(q.get('short_title_ratio',0)*100)}%",
        f"- Image/diagram slides: {int(q.get('image_slide_ratio',0)*100)}% — insert image-type slides at this rate",
        "",
        "Qualitative style:",
    ]
    for key, val in qual.items():
        if isinstance(val, list):
            val = "; ".join(str(v) for v in val)
        lines.append(f"- {key}: {val}")
    lines.append("")
    lines.append("Critical: build the deck as a narrative of many short slides, NOT a few dense slides. Mirror the titling voice, line length, and image rhythm. Do NOT produce mechanical labels like 'Topic — key points'.")
    block = "\n".join(lines)

    # Append user-customized dimension overrides if the teacher dragged sliders
    user_dims = profile.get("user_dimensions") or {}
    if user_dims:
        block += _user_dimensions_to_prompt_block(user_dims)
    return block
